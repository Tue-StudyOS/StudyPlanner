"""Build the StudyOS course presentation (English, dark, real app chrome).

The top bar is NOT rebuilt with shapes: it is the real thing, rendered from the
production markup/tokens via scripts/render-app-chrome.mjs (Playwright, 3x) and
placed as an image. One PNG per section, so the active pill moves while flipping.
The /log slide uses a faithful rendered example page the same way. Tool logos
come from scripts/render-tool-logos.mjs.

Flow (~15 min): nav sections mirror the grading rubric. Every rubric opens with
a divider (stop) slide, every content heading is "Section - Title" (plain
hyphen, one weight). No muted italic asides; those points are spoken. Footer on
framed slides, slide number in a small brand box, no shadows anywhere.

Output: docs/StudyOS-presentation.pptx
Rebuild: node scripts/render-app-chrome.mjs && python scripts/build-studyos-presentation.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SHOTS = ROOT / "docs" / "presentation-iteration-shots" / "dark-uniform"
CHROME = ROOT / "docs" / "app-chrome"
LOGOS = ROOT / "docs" / "tool-logos"
DIAGRAM = ROOT / "docs" / "StudyOS-architecture-diagram.png"
CLOUDFLARE = ROOT / "docs" / "cloudflare-logo.png"
ERSTIHEFT = ROOT / "docs" / "erstiheft-crop.png"
GPT_SHOT = ROOT / "docs" / "studyos-gpt-screenshot.png"
LOGO = ROOT / "docs" / "studyos-logo-glyph.png"
QR = ROOT / "docs" / "StudyOS-qr-site.png"
OUT = ROOT / "docs" / "StudyOS-presentation.pptx"

APP_URL = "https://studyplaner.pages.dev"
GPT_URL = "https://chatgpt.com/g/g-6a2de082a0b88191b833f7307d0c9429-studyos-bot"

BG = RGBColor(0x23, 0x26, 0x2D)
SURFACE = RGBColor(0x2E, 0x32, 0x3A)
BORDER = RGBColor(0x3C, 0x40, 0x4A)
FG = RGBColor(0xEC, 0xEC, 0xEC)
FG_MID = RGBColor(0xB4, 0xB6, 0xBC)
FG_MUTED = RGBColor(0x8A, 0x8D, 0x94)
LIGHT = RGBColor(0xF5, 0xF4, 0xF2)
PRIMARY = RGBColor(0xD1, 0x40, 0x60)  # app dark-mode primary (kept for accents if needed)

FOOTER = "StudyOS Practical - P. Gehler - July 2026"

W, H = Inches(13.333), Inches(7.5)
TOPBAR_H = 13.333 * 60 / 1440  # real bar aspect: 1440x60
FONT = "Inter"  # same as frontend --font-sans (Google Slides resolves it; local PPT may substitute)


def no_shadow(shape) -> None:
    """Strip every drop-shadow: inherit=False is not enough, PowerPoint still draws theme shadows."""
    try:
        shape.shadow.inherit = False
    except Exception:
        pass
    sp_pr = getattr(shape._element, "spPr", None)
    if sp_pr is None:
        return
    for tag in ("a:effectLst", "a:effectDag"):
        node = sp_pr.find(qn(tag))
        if node is not None:
            sp_pr.remove(node)


def text(slide, left, top, width, height, s: str, size: float, color: RGBColor, *,
         bold: bool = False, italic: bool = False, align=PP_ALIGN.LEFT,
         rotation: float = 0.0, link: str | None = None, line_gap: float | None = None):
    """Text box; newlines become real paragraphs (runs swallow them otherwise)."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(s.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_gap is not None:
            p.space_after = Pt(line_gap)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = FONT
    if link:
        # click action on the shape keeps our colors (no blue hyperlink styling)
        box.click_action.hyperlink.address = link
    if rotation:
        box.rotation = rotation
    return box


def fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    no_shadow(shape)


def card(slide, left, top, width, height, color: RGBColor = SURFACE):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill(c, color)
    c.line.color.rgb = BORDER
    c.line.width = Pt(1)
    try:
        c.adjustments[0] = 0.06
    except Exception:
        pass
    return c


def bg_fill(slide) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    fill(shape, BG)
    tree = slide.shapes._spTree
    el = shape._element
    tree.remove(el)
    tree.insert(2, el)


def app_frame(slide, active: str | None) -> None:
    """Real rendered top bar as an image, one variant per active section."""
    bg_fill(slide)
    name = active.lower().replace(" ", "-") if active else "none"
    bar = CHROME / f"topbar-{name}.png"
    if bar.exists():
        slide.shapes.add_picture(str(bar), 0, 0, width=W, height=Inches(TOPBAR_H))


def page_header(slide, section: str, title: str) -> None:
    """One heading line, one weight: 'Section - Title'."""
    text(slide, Inches(0.7), Inches(0.9), Inches(11.9), Inches(0.6),
         f"{section} - {title}", 26, FG, bold=True)


def add_image_fit(slide, path: Path, left_in: float, top_in: float, max_w_in: float, max_h_in: float,
                  rotation: float = 0.0, link: str | None = None):
    from PIL import Image
    if not path.exists():
        text(slide, Inches(left_in), Inches(top_in + max_h_in / 2), Inches(max_w_in), Inches(0.4),
             path.name, 12, FG_MUTED, align=PP_ALIGN.CENTER)
        return None
    with Image.open(path) as im:
        iw, ih = im.size
    aspect = iw / ih
    if aspect > max_w_in / max_h_in:
        w, h = max_w_in, max_w_in / aspect
    else:
        h, w = max_h_in, max_h_in * aspect
    pic = slide.shapes.add_picture(str(path), Inches(left_in + (max_w_in - w) / 2),
                                   Inches(top_in + (max_h_in - h) / 2), width=Inches(w), height=Inches(h))
    no_shadow(pic)
    if rotation:
        pic.rotation = rotation
    if link:
        pic.click_action.hyperlink.address = link
    return pic


def bullet_rows(slide, items: list[str], y: float, size: float = 17, step: float = 0.72,
                width: float = 11.2) -> None:
    for it in items:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), Inches(y + 0.13), Inches(0.09), Inches(0.09))
        fill(dot, FG_MID)
        text(slide, Inches(1.15), Inches(y), Inches(width), Inches(0.6), it, size, FG)
        y += step


def quote_grid(slide, quotes: list[str], top: float, size: float = 15) -> None:
    """Two calm columns of thought fragments: aligned rows, tiny alternating tilt."""
    col_x = (1.1, 7.2)
    col_w = 5.2
    row_h = (7.0 - top) / ((len(quotes) + 1) // 2)
    for i, q in enumerate(quotes):
        col, row = i % 2, i // 2
        x = col_x[col] + (0.25 if row % 2 else 0.0)
        rot = (-1.2, 1.2)[(i + row) % 2]
        text(slide, Inches(x), Inches(top + row * row_h), Inches(col_w), Inches(0.8),
             f"\u201c{q}\u201d", size, FG, italic=True, rotation=rot)


def footer(slide, number: int) -> None:
    """Gehler-style footer: deck context left, plain white slide number bottom right."""
    text(slide, Inches(0.35), Inches(7.13), Inches(8.0), Inches(0.3), FOOTER, 10, FG_MUTED)
    text(slide, Inches(12.55), Inches(7.1), Inches(0.6), Inches(0.35), str(number), 12, FG,
         align=PP_ALIGN.RIGHT)


def build() -> None:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]
    framed: set[int] = set()  # slides that get footer + number (everything but full-bleed shots)

    def s(active: str | None):
        slide = prs.slides.add_slide(blank)
        app_frame(slide, active)
        framed.add(id(slide))
        return slide

    def shot(fname: str) -> None:
        slide = prs.slides.add_slide(blank)
        bg_fill(slide)
        add_image_fit(slide, SHOTS / fname, 0.0, 0.0, 13.333, 7.5)

    def divider(section: str, rubric: str):
        """Stop slide at the start of each rubric: clear topic change, room to talk."""
        slide = s(section)
        text(slide, Inches(1.4), Inches(2.9), Inches(10.5), Inches(1.0), section, 44, FG,
             bold=True, align=PP_ALIGN.CENTER)
        text(slide, Inches(1.4), Inches(4.05), Inches(10.5), Inches(0.5), rubric, 15, FG_MUTED,
             align=PP_ALIGN.CENTER)
        return slide

    # ---- Title: the whole team ------------------------------------------------
    slide = s(None)
    text(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.2), "StudyOS", 54, FG, bold=True)
    text(slide, Inches(0.95), Inches(3.65), Inches(11), Inches(0.5), "Plan your studies in one place", 20, FG_MID)
    text(slide, Inches(0.95), Inches(4.9), Inches(11.5), Inches(0.5),
         "Ben Tischberger · Yonatan Dankner · Emre Sözbilir · Lena Binder", 15, FG_MUTED)
    text(slide, Inches(0.95), Inches(5.45), Inches(11), Inches(0.4), "Presented by Ben and Yonatan", 12, FG_MUTED)

    # ---- Problem ---------------------------------------------------------------
    divider("Problem", "problem framing and evolution")

    slide = s("Problem")
    page_header(slide, "Problem", "What we solve")
    bullet_rows(slide, [
        "Course info lives in ALMA, PDFs and spreadsheets",
        "Timetables need manual cross-checking",
        "Degree progress is guesswork",
    ], 2.05, size=17, step=0.55)
    # the course's X / Y / Z framing, stated explicitly
    labels = [("X", "Course planning is scattered across systems that don't talk to each other"),
              ("Y", "Informatik students in Tübingen, all six PO 2021 programs"),
              ("Z", "One app: catalog, weekly timetable and degree progress together")]
    card(slide, Inches(0.7), Inches(3.9), Inches(11.9), Inches(2.7))
    ly = 4.2
    for label, body in labels:
        text(slide, Inches(1.05), Inches(ly), Inches(0.9), Inches(0.45), label, 15, FG, bold=True)
        text(slide, Inches(2.0), Inches(ly), Inches(10.2), Inches(0.5), body, 14, FG)
        ly += 0.7

    # ---- Students ----------------------------------------------------------------
    divider("Students", "user research and evidence")

    slide = s("Students")
    page_header(slide, "Students", "What they told us")
    quote_grid(slide, [
        "ALMA is my only info source, but I hate the UI.",
        "I plan with Excel and the PO PDF.",
        "I check three sites before every semester.",
        "Show me which courses fit my degree.",
        "Let me search by professor.",
        "One calendar with all my courses.",
        "Tell me what I still need to graduate.",
        "Just make it one page.",
    ], 2.3)

    # ---- Product -----------------------------------------------------------------
    slide = divider("Product", "concept and solution design")
    add_image_fit(slide, LOGO, 6.24, 1.85, 0.85, 0.85)

    slide = s("Product")
    page_header(slide, "Product", "What StudyOS can do")
    cats = [
        ("Catalog", "All Informatik courses\nFilters and professor search\nNine semesters of data"),
        ("Planner", "Weekly timetable\nConflict check and tutorial groups\nCalendar export (ICS)"),
        ("Progress", "Transcript PDF upload\nExam regulations tracked\nSee what is still open"),
        ("Beyond the app", "GPT access via MCP\nFeedback built in\nGerman and English"),
    ]
    positions = [(0.7, 1.95), (6.9, 1.95), (0.7, 4.55), (6.9, 4.55)]
    for (title, body), (x, y) in zip(cats, positions):
        card(slide, Inches(x), Inches(y), Inches(5.7), Inches(2.4))
        text(slide, Inches(x + 0.35), Inches(y + 0.25), Inches(5.0), Inches(0.45), title, 18, FG, bold=True)
        text(slide, Inches(x + 0.35), Inches(y + 0.8), Inches(5.0), Inches(1.5), body, 13.5, FG_MID, line_gap=4)

    slide = s("Product")
    page_header(slide, "Product", "Degree programs")
    text(slide, Inches(0.72), Inches(1.55), Inches(11.9), Inches(0.45),
         "We support the PO 2021 exam regulations of six programs. Everything works for each of them.",
         15, FG)
    card(slide, Inches(0.7), Inches(2.35), Inches(5.9), Inches(4.35))
    text(slide, Inches(1.05), Inches(2.65), Inches(5.2), Inches(0.45), "Bachelor (180 ECTS)", 17, FG, bold=True)
    text(slide, Inches(1.05), Inches(3.4), Inches(5.2), Inches(3.0),
         "B.Sc. Informatik\nB.Sc. Bioinformatik\nB.Sc. Medieninformatik\nB.Sc. Medizininformatik",
         15, FG_MID, line_gap=10)
    card(slide, Inches(6.9), Inches(2.35), Inches(5.9), Inches(4.35))
    text(slide, Inches(7.25), Inches(2.65), Inches(5.2), Inches(0.45), "Master (120 ECTS)", 17, FG, bold=True)
    text(slide, Inches(7.25), Inches(3.4), Inches(5.2), Inches(3.0),
         "M.Sc. Informatik\nM.Sc. Machine Learning",
         15, FG_MID, line_gap=10)

    # ---- Build -------------------------------------------------------------------
    divider("Build", "implementation and prototype")

    slide = s("Build")
    page_header(slide, "Build", "How the data flows")
    add_image_fit(slide, DIAGRAM, 0.35, 1.55, 12.6, 5.5)

    slide = s("Build")
    page_header(slide, "Build", "Runs on Cloudflare")
    bullet_rows(slide, [
        "Pages serves the frontend, Workers run the API",
        "D1: SQLite at the edge, one database for catalog and accounts",
        "Merge to main deploys automatically, config checked in CI",
        "Free tier covers all of it",
    ], 2.3, size=16, width=7.2)
    card(slide, Inches(8.8), Inches(2.5), Inches(3.8), Inches(2.0), LIGHT)
    add_image_fit(slide, CLOUDFLARE, 9.1, 2.75, 3.2, 1.5)

    slide = s("Build")
    page_header(slide, "Build", "How we worked")
    # three cards: Claude, Cursor, Codex+Pi together
    cards = [
        ("Claude", "Great at UI work. Expensive.", [("claude.png", 1.15)]),
        ("Cursor", "Cheap, solid code. UI taste is meh.", [("cursor.png", 1.15)]),
        ("Codex + Pi", "Both used with GPT.",
         [("codex.png", 0.95), ("pi.png", 0.95)]),
    ]
    cw, gap = 3.9, 0.25
    cx = (13.333 - (cw * 3 + gap * 2)) / 2
    for i, (title, body, logos) in enumerate(cards):
        x = cx + i * (cw + gap)
        card(slide, Inches(x), Inches(2.05), Inches(cw), Inches(3.15))
        logo_w = 1.05
        total = logo_w * len(logos) + 0.2 * (len(logos) - 1)
        lx = x + (cw - total) / 2
        for logo, lh in logos:
            add_image_fit(slide, LOGOS / logo, lx, 2.3, logo_w, lh)
            lx += logo_w + 0.2
        text(slide, Inches(x + 0.25), Inches(3.55), Inches(cw - 0.5), Inches(0.4), title, 16, FG,
             bold=True, align=PP_ALIGN.CENTER)
        text(slide, Inches(x + 0.25), Inches(4.05), Inches(cw - 0.5), Inches(1.0), body, 12.5, FG_MID,
             align=PP_ALIGN.CENTER, line_gap=3)
    # takeaway: quiet text arrow, same weight as the sentence
    text(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.5),
         "\u2192  Each tool has clear limits. We mixed them by task and reviewed everything.",
         16, FG, align=PP_ALIGN.CENTER)

    slide = s("Build")
    page_header(slide, "Build", "StudyOS Bot, our GPT")
    card(slide, Inches(0.7), Inches(2.1), Inches(7.3), Inches(4.2))
    add_image_fit(slide, GPT_SHOT, 0.95, 2.35, 6.8, 3.7, link=GPT_URL)
    ry = 2.45
    for t in [
        "Search courses, resolve one, get details",
        "Works with ChatGPT, Claude, any MCP client",
        "Endpoint: studyplaner.pages.dev/mcp",
    ]:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.35), Inches(ry + 0.13), Inches(0.08), Inches(0.08))
        fill(dot, FG_MID)
        text(slide, Inches(8.6), Inches(ry), Inches(4.2), Inches(0.8), t, 14, FG)
        ry += 0.85
    text(slide, Inches(8.6), Inches(5.3), Inches(4.2), Inches(0.4), "Open the StudyOS Bot", 14, FG, bold=True, link=GPT_URL)

    # ---- Iterate -------------------------------------------------------------------
    slide = divider("Iterate", "testing, iteration and shipping")
    # double diamond: plan (discover/define), then implement (develop/deliver)
    for dx in (5.17, 6.67):
        d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(dx), Inches(4.75), Inches(1.5), Inches(1.5))
        d.fill.background()
        d.line.color.rgb = FG_MID
        d.line.width = Pt(2)
        no_shadow(d)

    # Fast click-through, whole original screenshots, no chrome, no labels.
    # Compass v1 dropped; one feedback stop in the middle; beta before prod.
    for fname in [
        "02-compass-v2.png", "03-html-dashboard.png", "04-html-catalog.png",
        "05-early-react-catalog.png", "06-early-catalog-cards.png", "07-api-catalog.png",
        "08-studyplanner-catalog.png", "09-filters-catalog.png", "10-test-landing.png",
        "11-test-catalog.png",
    ]:
        shot(fname)

    slide = s("Iterate")
    page_header(slide, "Iterate", "Feedback along the way")
    quote_grid(slide, [
        "Icons for the semesters instead of writing them out.",
        "Make more things clickable, let them set the filters.",
        "Let me drag courses in and out of the calendar.",
        "Delete looks like collapse, I lost my plan.",
        "Why is everything red? Red means error.",
        "Two lecture halls, but no real conflict.",
    ], 2.3)

    for fname in ["12-season-catalog.png", "13-currentish-catalog.png", "15-prod-beta.png", "14-prod-catalog.png"]:
        shot(fname)

    slide = s("Iterate")
    text(slide, Inches(1.4), Inches(3.3), Inches(10.5), Inches(0.8),
         "…and many more changes than fit in one deck.", 26, FG, bold=True, align=PP_ALIGN.CENTER)

    slide = s("Iterate")
    page_header(slide, "Iterate", "What we fixed")
    pairs = [
        ("Favorites in terms with no offering", "Availability check"),
        ("Fake room conflicts", "Fixed conflict logic"),
        ("Delete looked like collapse", "Clearer delete"),
        ("Too much red", "Red only for errors"),
        ("Crowded planner", "Decluttered layout"),
    ]
    y = 2.25
    for heard, fixed in pairs:
        text(slide, Inches(0.9), Inches(y), Inches(5.6), Inches(0.5), heard, 16, FG_MID)
        text(slide, Inches(6.5), Inches(y), Inches(0.4), Inches(0.5), "\u2192", 16, FG, bold=True)
        text(slide, Inches(7.1), Inches(y), Inches(5.5), Inches(0.5), fixed, 16, FG)
        y += 0.7

    slide = s("Iterate")
    page_header(slide, "Iterate", "Where we are today")
    bullet_rows(slide, [
        "1,158 courses across 9 semesters, summer 2022 to summer 2026",
        "Live at studyplaner.pages.dev, catalog public, account for plans and progress",
        "Timetable with ICS export, transcript import, German and English",
    ], 2.3, size=17)

    # ---- Demo: just the link, clickable ----------------------------------------
    slide = s("Demo")
    text(slide, Inches(0.9), Inches(3.2), Inches(11.5), Inches(1.0), "Live demo", 36, FG, bold=True, link=APP_URL)

    # ---- Future --------------------------------------------------------------------
    divider("Future", "ideas, maintenance, what could ship next")

    # Log dominates the slide; no Codex logo so the screenshot can use the space
    slide = s("Future")
    page_header(slide, "Future", "Watch the log, ship the fix")
    bullet_rows(slide, [
        "The request log already captures failures",
        "Agents can watch it and propose fixes",
        "Scraping and deploys already run on their own",
        "Missing: PO knowledge beyond our programs",
    ], 2.05, size=14, step=0.58, width=3.6)
    add_image_fit(slide, CHROME / "log-example-cropped.png", 4.4, 1.35, 8.6, 5.7)

    # Fachschaft / distribution, separate slide
    slide = s("Future")
    page_header(slide, "Future", "Reach and handoff")
    bullet_rows(slide, [
        "Our hope: the Fachschaft takes it over, we are in contact",
        "Erstiheft and the CS mailing list to reach new students",
        "Maintenance stays light if scraping and deploys keep running",
    ], 2.3, size=16, width=7.4)
    add_image_fit(slide, ERSTIHEFT, 9.1, 1.7, 3.5, 5.2)

    # ---- Try it: one link left, clickable QR right, last slide ------------------
    slide = s("Try it")
    page_header(slide, "Try it", "Scan or click")
    text(slide, Inches(0.9), Inches(2.5), Inches(6.2), Inches(0.6), "studyplaner.pages.dev", 24, FG, bold=True, link=APP_URL)
    text(slide, Inches(0.9), Inches(3.4), Inches(5.8), Inches(1.5),
         "The feedback button is in the app. Use it, we read everything.", 16, FG_MID)
    card(slide, Inches(7.6), Inches(1.9), Inches(4.9), Inches(5.0))
    add_image_fit(slide, QR, 8.35, 2.7, 3.4, 3.4, link=APP_URL)

    # Footer + slide number on every framed slide (full-bleed shots stay clean)
    for i, slide in enumerate(prs.slides, start=1):
        if id(slide) in framed:
            footer(slide, i)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT} ({OUT.stat().st_size} bytes, {len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
